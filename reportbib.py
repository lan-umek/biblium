import os
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font
from openpyxl.utils import get_column_letter
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.worksheet.hyperlink import Hyperlink
from openpyxl.drawing.image import Image as XLImage
from utilsbib import to_excel_fancy

COLOR_MAP = {
    "LightGreen": "90EE90",
    "LightBlue": "ADD8E6",
    "Khaki": "F0E68C",
    "Gold": "FFD700",
    "Orange": "FFA500",
    "Pink": "FFC0CB",
    "Red": "FF0000",
    "Blue": "0000FF",
    "Green": "008000",
    "Purple": "800080",
    "Yellow": "FFFF00",
    "Gray": "808080",
    "Cyan": "00FFFF",
    "Magenta": "FF00FF",
    "White": "FFFFFF",
    "Black": "000000"
}

def resolve_color(name, default="FFFFFF"):
    return COLOR_MAP.get(name, COLOR_MAP.get(default))

def insert_icon(toc_ws, fd, row_idx, icon_type="table"):
    icon_filename = "table.png" if icon_type == "table" else "plot.png"
    icon_path = os.path.join(fd, "additional files", "icons", icon_filename)
    if os.path.exists(icon_path):
        img = XLImage(icon_path)
        img.width = img.height = 20
        toc_ws.add_image(img, f"D{row_idx}")

fd = os.path.dirname(__file__)

def save_excel_report_from_template(
    self,
    output_path: str = "bibliometric_report.xlsx",
    template_path: str = fd + "\\additional files\\template for excel output.xlsx",
    template_sheet: str = "testing",
    striped_rows: bool = False,
    show_icons: bool = False,
    autofit: bool = True,
    show_gridlines: bool = True,
    freeze_header: bool = True,
    **fancy_kwargs
):
    """
    Generate and save an Excel report based on a template sheet,
    adding tables and plots defined in the template only if their
    corresponding data attributes or plot files exist.

    Parameters:
        output_path (str): Path to save the generated Excel workbook.
        template_path (str): Path to the Excel template file.
        template_sheet (str): Sheet name within the template file.
        striped_rows (bool): Whether to apply alternating row colors in the TOC.
        show_icons (bool): Whether to insert icons for tables/plots in the TOC.
        autofit (bool): Whether to auto-adjust column widths for included tables.
        show_gridlines (bool): Whether to display gridlines in all sheets.
        freeze_header (bool): Whether to freeze the header row in detail sheets.
        **fancy_kwargs: Additional styling or features (currently unused).
    """
    # Read the template definitions
    df_template = pd.read_excel(template_path, sheet_name=template_sheet)

    wb = Workbook()
    toc_ws = wb.active
    toc_ws.title = "Table of Contents"
    toc_ws.freeze_panes = "A3"
    toc_tab_color = resolve_color("LightBlue")
    if toc_tab_color:
        toc_ws.sheet_properties.tabColor = toc_tab_color

    # TOC header
    header = ["Level 1", "Level 2", "Sheet Name"] + (["Icon"] if show_icons else [])
    toc_ws.append(header)
    for col_num in range(1, len(header) + 1):
        cell = toc_ws.cell(row=1, column=col_num)
        cell.fill = PatternFill(start_color="D9EAD3", end_color="D9EAD3", fill_type="solid")
        cell.font = Font(bold=True)

    row_idx = 2
    for _, row in df_template.iterrows():
        level1 = str(row.get("Level 1", "")).strip()
        level2 = str(row.get("Level 2", "")).strip()
        sheet_name = f"{level2}"[:31]

        include_table = str(row.get("Include Table", "")).strip().upper() == "TRUE"
        data_attr = str(row.get("Data Attr", "")).strip()
        df_obj = getattr(self, data_attr, None)
        table_exists = isinstance(df_obj, pd.DataFrame)

        include_plot = str(row.get("Include Plot", "")).strip().upper() == "TRUE"
        path = str(row.get("Path", "")).strip()
        plot_filename = str(row.get("Plot Filename", "")).strip()
        full_plot_path = os.path.abspath(os.path.join(fd, path, plot_filename)) if plot_filename else None
        plot_exists = include_plot and full_plot_path and os.path.exists(full_plot_path)

        # Skip entries without data or plot
        if not table_exists and not plot_exists:
            continue

        # Determine tab and toc colors
        tab_color = resolve_color(str(row.get("Tab Color", "")).strip())
        toc_color = resolve_color(str(row.get("TOC Label Color", "")).strip())

        # Create sheet
        ws = wb.create_sheet(sheet_name)
        if tab_color:
            ws.sheet_properties.tabColor = tab_color

        # Populate TOC
        toc_ws.cell(row=row_idx, column=1, value=level1)
        toc_ws.cell(row=row_idx, column=2, value=level2)
        toc_link = toc_ws.cell(row=row_idx, column=3, value=sheet_name)
        toc_link.hyperlink = Hyperlink(ref=toc_link.coordinate, location=f"'{sheet_name}'!A1")
        toc_link.font = Font(color="0000FF", underline="single")
        if toc_color:
            toc_link.fill = PatternFill(start_color=toc_color, end_color=toc_color, fill_type="solid")

        if show_icons:
            icon_type = "table" if table_exists else "plot"
            insert_icon(toc_ws, fd, row_idx, icon_type)

        row_idx += 1

        # Insert table if exists
        if table_exists:
            float_cols = [
                col for col in df_obj.columns
                if pd.api.types.is_float_dtype(df_obj[col])
                or (pd.api.types.is_numeric_dtype(df_obj[col]) and not pd.api.types.is_integer_dtype(df_obj[col]))
            ]
            if freeze_header:
                ws.freeze_panes = "A2"
            for r_idx, row_data in enumerate(dataframe_to_rows(df_obj, index=False, header=True), 1):
                for c_idx, value in enumerate(row_data, 1):
                    if c_idx - 1 < len(df_obj.columns):
                        colname = df_obj.columns[c_idx - 1]
                        if colname in float_cols and isinstance(value, float):
                            value = round(value, 3)
                    cell = ws.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == 1:
                        cell.font = Font(bold=True)

        # Insert plot if exists
        if plot_exists:
            try:
                img = XLImage(full_plot_path)
                img.width = img.height = 0.4 * img.width
                ws.add_image(img, "C4")
                ws["A1"] = "Path to the plot (for reference only)"
                ws["A2"] = full_plot_path
            except Exception:
                pass

    # Style TOC
    if striped_rows:
        for row in toc_ws.iter_rows(min_row=2, max_row=row_idx - 1):
            if row[0].row % 2 == 0:
                for cell in row:
                    cell.fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")

    for cols in toc_ws.columns:
        width = max(len(str(cell.value)) if cell.value else 0 for cell in cols) + 2
        toc_ws.column_dimensions[get_column_letter(cols[0].column)].width = width

    if "Sheet" in wb.sheetnames:
        del wb["Sheet"]

    if autofit:
        for sheet in wb.worksheets:
            for cols in sheet.columns:
                width = max(len(str(cell.value)) if cell.value else 0 for cell in cols) + 2
                sheet.column_dimensions[get_column_letter(cols[0].column)].width = width

    if not show_gridlines:
        for sheet in wb.worksheets:
            sheet.sheet_view.showGridLines = False

    wb.save(output_path)
    

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import pandas as pd

def save_powerpoint_report_from_template(
    self,
    output_path: str = "bibliometric_report.pptx",
    template_path: str = fd + "\\additional files\\template for powerpoint output.xlsx",
    template_sheet: str = "testing",
    theme_color: str = "FFFFFF",
    title_font: dict = {"name": "Arial", "size": 36, "bold": True},
    subtitle_font: dict = {"name": "Arial", "size": 24, "italic": True},
    bullet_font: dict = {"name": "Calibri", "size": 18},
    image_width: float = 6.5,
    align_center: bool = True,
    include_toc_slide: bool = True,
    footer_text: str = "",
    top_n: int = 5,
    presentation_title: str = "Bibliometric Report",
    presentation_subtitle: str = "",
    bold_slide_titles: bool = True,
    custom_template_path: str = "",
    color_by_section: bool = False,
    style_table_headers: bool = False,
    show_section_footer: bool = False,
    insert_logo: bool = False,
    logo_path: str = fd + "\\additional files\\icons\\logo.png",
    logo_position: str = "bottom-right",
    logo_width: float = 1.0,
    theme_palette: dict = {
        "header_fill": RGBColor(220, 220, 220),
        "header_font": RGBColor(0, 0, 0),
        "cell_border": RGBColor(200, 200, 200)
    }
):
    prs = Presentation(custom_template_path) if custom_template_path else Presentation()
        
    insert_logo = True
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_shape = title_slide.shapes.title
    title_shape.text = presentation_title
    title_shape.text_frame.paragraphs[0].font.name = title_font.get("name", "Arial")
    title_shape.text_frame.paragraphs[0].font.size = Pt(title_font.get("size", 36))
    title_shape.text_frame.paragraphs[0].font.bold = title_font.get("bold", True)

    if presentation_subtitle:
        subtitle_shape = title_slide.placeholders[1]
        subtitle_shape.text = presentation_subtitle
        subtitle_shape.text_frame.paragraphs[0].font.name = subtitle_font.get("name", "Arial")
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(subtitle_font.get("size", 24))
        subtitle_shape.text_frame.paragraphs[0].font.italic = subtitle_font.get("italic", False)

    df_template = pd.read_excel(template_path, sheet_name=template_sheet)
    bullet_slide_layout = prs.slide_layouts[1]

    if include_toc_slide:
        toc_slide = prs.slides.add_slide(bullet_slide_layout)
        toc_slide.shapes.title.text = "Table of Contents"
        if bold_slide_titles:
            toc_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        toc_tf = toc_slide.placeholders[1].text_frame
        level1_groups = df_template['Level 1'].dropna().unique()
        for lvl1 in level1_groups:
            p = toc_tf.add_paragraph()
            p.text = lvl1
            p.level = 0

        for lvl1 in level1_groups:
            section_color = RGBColor(230, 230, 250) if color_by_section else None
            section_slide = prs.slides.add_slide(bullet_slide_layout)
            section_slide.shapes.title.text = lvl1
            if bold_slide_titles:
                section_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            if color_by_section:
                fill = section_slide.background.fill
                fill.solid()
                fill.fore_color.rgb = section_color
            tf1 = section_slide.placeholders[1].text_frame
            level2_items = df_template[df_template['Level 1'] == lvl1]['Level 2'].dropna().unique()
            for lvl2 in level2_items:
                p = tf1.add_paragraph()
                p.text = lvl2
                p.level = 0

            for _, row in df_template[df_template['Level 1'] == lvl1].iterrows():
                level2 = str(row.get("Level 2", "")).strip()
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = level2
                if bold_slide_titles:
                    title_shape.text_frame.paragraphs[0].font.bold = True
                placeholder = slide.placeholders[1]

                include_plot = str(row.get("Include Plot", "")).strip().upper() == "TRUE"
                path = str(row.get("Path", "")).strip()
                filename = str(row.get("Plot Filename", "")).strip()
                if include_plot and filename:
                    img_path = os.path.join(path, filename)
                    if os.path.exists(img_path):
                        left = placeholder.left
                        top = placeholder.top
                        width = placeholder.width
                        height = placeholder.height
                        
                        
                        slide.shapes.add_picture(img_path, left, top, width=width, height=height)

                include_table = str(row.get("Include Table", "")).strip().upper() == "TRUE"
                data_attr = str(row.get("Data Attr", "")).strip()
                if include_table and hasattr(self, data_attr):
                    df = getattr(self, data_attr)
                    if isinstance(df, pd.DataFrame) and not df.empty:
                        df_preview = df.head(top_n).copy()
                        for col in df_preview.columns:
                            if pd.api.types.is_float_dtype(df_preview[col]):
                                df_preview[col] = df_preview[col].round(3)
                        rows, cols = df_preview.shape
                        left = placeholder.left
                        top = placeholder.top
                        width = placeholder.width
                        height = placeholder.height
                        table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
                        for j, col_name in enumerate(df_preview.columns):
                            cell = table_shape.cell(0, j)
                            cell.text = str(col_name)
                            if style_table_headers:
                                cell.text_frame.paragraphs[0].font.bold = True
                                cell.text_frame.paragraphs[0].font.color.rgb = theme_palette["header_font"]
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = theme_palette["header_fill"]
                                
                        for i in range(rows):
                            for j in range(cols):
                                cell = table_shape.cell(i + 1, j)
                                cell.text = str(df_preview.iloc[i, j])
                                

                narrative_attr = str(row.get("Narrative", "")).strip()
                if hasattr(self, narrative_attr):
                    narrative = getattr(self, narrative_attr)
                    if narrative:
                        textbox = slide.shapes.add_textbox(placeholder.left, Inches(5), Inches(8), Inches(2))
                        tf = textbox.text_frame
                        for line in narrative.splitlines():
                            p = tf.add_paragraph()
                            p.text = line
                            p.font.size = Pt(bullet_font.get("size", 18))
                            p.font.name = bullet_font.get("name", "Calibri")

                if show_section_footer:
                    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(8), Inches(0.3))
                    footer.text_frame.text = lvl1

    print(f"Checking logo path: {logo_path}")
    if insert_logo and os.path.exists(logo_path):
        print(f"Logo found, inserting from: {logo_path}")
        for slide in prs.slides:
            if logo_position == "bottom-right":
                left, top = Inches(9), Inches(6.5)
            elif logo_position == "bottom-left":
                left, top = Inches(0.5), Inches(6.5)
            elif logo_position == "top-left":
                left, top = Inches(0.5), Inches(0.5)
            elif logo_position == "top-right":
                left, top = Inches(9), Inches(0.5)
            else:
                left, top = Inches(9), Inches(6.5)
            slide.shapes.add_picture(logo_path, left, top, width=Inches(logo_width))

    prs.save(output_path)
    print(f"PowerPoint saved to {output_path}")
    
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.style import WD_STYLE_TYPE
import os
import pandas as pd

def save_word_report_from_template(
    self,
    output_path: str = "bibliometric_report.docx",
    template_path: str = fd + "\\additional files\\template for word output.xlsx",
    template_sheet: str = "testing",
    top_n: int = 50,
    default_table_style: str = "Light List Accent 1",
    enumerate_figures: bool = True,
    enumerate_tables: bool = True,
    enumerate_sections: bool = True,
    include_logo: bool = True,
    heading_color: str = "blue"):
    """
    Generate a Word report based on Excel-configured structure.
    Inserts tables, plots, captions, and text blocks based on user selections.
    """
    doc = Document()
    
    from docx.shared import RGBColor
    color_map = {
        "orange": RGBColor(255, 165, 0),
        "blue": RGBColor(0, 112, 192),
        "black": RGBColor(0, 0, 0),
        "gray": RGBColor(128, 128, 128),
        "green": RGBColor(0, 176, 80),
        "red": RGBColor(255, 0, 0),
        "purple": RGBColor(112, 48, 160)
    }
    title_para = doc.add_heading("Bibliometric Report", level=0)
    title_para.runs[0].font.color.rgb = color_map.get(heading_color.lower(), RGBColor(255, 165, 0))

    toc = []
    figure_counter = 1
    table_counter = 1

    df_template = pd.read_excel(template_path, sheet_name=template_sheet)
    df_template.columns = df_template.columns.str.strip().str.replace(" ", "_").str.lower()

    from docx.shared import RGBColor
    for i1, lvl1 in enumerate(df_template['level_1'].dropna().unique(), 1):
        heading1 = f"{i1}. {lvl1}" if enumerate_sections else str(lvl1)
        color_map = {
            "orange": RGBColor(255, 165, 0),
            "blue": RGBColor(0, 112, 192),
            "black": RGBColor(0, 0, 0),
            "gray": RGBColor(128, 128, 128),
            "green": RGBColor(0, 176, 80),
            "red": RGBColor(255, 0, 0),
            "purple": RGBColor(112, 48, 160)
        }
        para1 = doc.add_heading(heading1, level=1)
        run1 = para1.runs[0]
        run1.font.color.rgb = color_map.get(heading_color.lower(), RGBColor(255, 165, 0))
        toc.append((heading1, 1))

        for i2, (_, row) in enumerate(df_template[df_template['level_1'] == lvl1].iterrows(), 1):
            lvl2 = str(row.get("level_2", "")).strip()
            if lvl2:
                heading2 = f"{i1}.{i2} {lvl2}" if enumerate_sections else lvl2
                para2 = doc.add_heading(heading2, level=2)
                para2.runs[0].font.color.rgb = color_map.get(heading_color.lower(), RGBColor(255, 165, 0))
                toc.append((heading2, 2))

            # Insert optional description text
            desc_attr = str(row.get("description", "")).strip()
            if hasattr(self, desc_attr):
                desc_text = getattr(self, desc_attr)
                if isinstance(desc_text, str):
                    doc.add_paragraph(desc_text)

            # Insert narrative (optional)
            narrative_attr = str(row.get("narrative", "")).strip()
            if hasattr(self, narrative_attr):
                narrative = getattr(self, narrative_attr)
                if isinstance(narrative, str):
                    for line in narrative.splitlines():
                        doc.add_paragraph(line)

            # Table
            include_table = str(row.get("include_table", "")).strip().upper() == "TRUE"
            data_attr = str(row.get("data_attr", "")).strip()
            if include_table and hasattr(self, data_attr):
                df = getattr(self, data_attr)
                if isinstance(df, pd.DataFrame) and not df.empty:
                    df_preview = df.head(top_n).copy()
                    for col in df_preview.columns:
                        if pd.api.types.is_float_dtype(df_preview[col]):
                            df_preview[col] = df_preview[col].round(3)
                    rows, cols = df_preview.shape
                    table = doc.add_table(rows=rows + 1, cols=cols)
                    table.style = str(row.get("table_style") or default_table_style)
                    for j, col_name in enumerate(df_preview.columns):
                        table.cell(0, j).text = str(col_name)
                    for i in range(rows):
                        for j in range(cols):
                            table.cell(i + 1, j).text = str(df_preview.iloc[i, j])
                    # Add caption
                    caption = str(row.get("caption", "")).strip()
                    if caption:
                        caption_prefix = f"Table {table_counter}: " if enumerate_tables else ""
                        table_counter += 1
                        caption_para = doc.add_paragraph(caption_prefix + caption, style='Caption')
                        caption_para.runs[0].font.color.rgb = color_map.get(heading_color.lower(), RGBColor(255, 165, 0))

            # Plot
            include_plot = str(row.get("include_plot", "")).strip().upper() == "TRUE"
            path = str(row.get("path", "")).strip()
            filename = str(row.get("plot_filename", "")).strip()
            if include_plot and filename:
                img_path = os.path.join(path, filename)
                if os.path.exists(img_path):
                    doc.add_picture(img_path, width=Inches(5.5))
                    caption = str(row.get("caption", "")).strip()
                    if caption:
                        caption_prefix = f"Figure {figure_counter}: " if enumerate_figures else ""
                        figure_counter += 1
                        doc.add_paragraph(caption_prefix + caption, style='Caption')

            # Page break if needed
            page_break = str(row.get("page_break_after", "")).strip().upper() == "TRUE"
            if page_break:
                doc.add_page_break()


    # Move TOC + title section to the beginning
    paragraphs = []
    paragraphs.append(doc.add_heading("Bibliometric Report", level=0))
    toc_paragraph = doc.add_paragraph()
    run = toc_paragraph.add_run("Table of Contents")
    run.bold = True
    run.font.size = Pt(14)
    paragraphs.append(toc_paragraph)
    for text, level in toc:
        p = doc.add_paragraph(text)
        paragraphs.append(p)
    paragraphs.append(doc.add_paragraph())  # empty paragraph before content
    paragraphs.append(doc.add_page_break())  # page break after TOC
    for p in reversed(paragraphs):
        doc._body._element.remove(p._element)
        doc._body._element.insert(0, p._element)

    # Add logo to footer of each section
    if include_logo:
        for section in doc.sections:
            footer = section.footer
            if not footer.paragraphs:
                footer_paragraph = footer.add_paragraph()
            else:
                footer_paragraph = footer.paragraphs[0]
            logo_path = os.path.join(fd, "additional files", "icons", "logo.png")
            if os.path.exists(logo_path):
                run = footer_paragraph.add_run()
                run.add_picture(logo_path, width=Inches(1.0))

    doc.save(output_path)
    print(f"Word report saved to {output_path}")


def save_tex_report_from_template(
    self,
    output_path: str = "bibliometric_report.tex",
    template_path: str = fd + "\\additional files\\template for tex output.xlsx",
    template_sheet: str = "testing",
    top_n: int = 50,
    enumerate_figures: bool = True,
    enumerate_tables: bool = True,
    enumerate_sections: bool = True
):
    """
    Generate a LaTeX report (.tex) based on an Excel-configured structure.
    """
    import pandas as pd
    import os

    df_template = pd.read_excel(template_path, sheet_name=template_sheet)
    df_template.columns = df_template.columns.str.strip().str.replace(" ", "_").str.lower()

    figure_counter = 1
    table_counter = 1
    lines = []

    lines.append("\\documentclass{article}")
    lines.append("\\usepackage[utf8]{inputenc}")
    lines.append("\\usepackage{graphicx}")
    lines.append("\\usepackage{booktabs}")
    lines.append("\\usepackage{geometry}")
    lines.append("\\geometry{margin=1in}")
    lines.append("\\title{Bibliometric Report}")
    lines.append("\\date{}")
    lines.append("\\begin{document}")
    lines.append("\\maketitle")
    lines.append("")

    for i1, lvl1 in enumerate(df_template['level_1'].dropna().unique(), 1):
        heading1 = f"{i1}. {lvl1}" if enumerate_sections else lvl1
        lines.append(f"\\section{{{heading1}}}")

        for i2, (_, row) in enumerate(df_template[df_template['level_1'] == lvl1].iterrows(), 1):
            lvl2 = str(row.get("level_2", "")).strip()
            if lvl2:
                heading2 = f"{i1}.{i2} {lvl2}" if enumerate_sections else lvl2
                lines.append(f"\\subsection{{{heading2}}}")

            desc_attr = str(row.get("description", "")).strip()
            if hasattr(self, desc_attr):
                desc_text = getattr(self, desc_attr)
                if isinstance(desc_text, str):
                    lines.append(desc_text + "\\\n")

            narrative_attr = str(row.get("narrative", "")).strip()
            if hasattr(self, narrative_attr):
                narrative = getattr(self, narrative_attr)
                if isinstance(narrative, str):
                    for line in narrative.splitlines():
                        lines.append(line + "\\\n")

            include_table = str(row.get("include_table", "")).strip().upper() == "TRUE"
            data_attr = str(row.get("data_attr", "")).strip()
            if include_table and hasattr(self, data_attr):
                df = getattr(self, data_attr)
                if isinstance(df, pd.DataFrame) and not df.empty:
                    df_preview = df.head(top_n).copy()
                    latex_table = df_preview.to_latex(index=False, escape=True)
                    lines.append(latex_table)
                    caption = str(row.get("caption", "")).strip()
                    if caption:
                        caption_prefix = f"Table {table_counter}: " if enumerate_tables else ""
                        lines.append(f"\textbf{{{caption_prefix}{caption}}}\
")
                        table_counter += 1

            include_plot = str(row.get("include_plot", "")).strip().upper() == "TRUE"
            path = str(row.get("path", "")).strip()
            filename = str(row.get("plot_filename", "")).strip()
            if include_plot and filename:
                img_path = os.path.join(path, filename)
                if os.path.exists(img_path):
                    # normalize path with forward slashes and escape underscores for LaTeX
                    tex_img_path = img_path.replace("\\\\", "/").replace("_", "\\_")
                    lines.append(f"\begin{{center}}\includegraphics[width=0.9\textwidth]{{{tex_img_path}}}\end{{center}}")
                    caption = str(row.get("caption", "")).strip()
                    if caption:
                        caption_prefix = f"Figure {figure_counter}: " if enumerate_figures else ""
                        lines.append(f"\textit{{{caption_prefix}{caption}}}\\")
                    figure_counter += 1

    lines.append("\\end{document}")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"TeX report saved to {output_path}")
